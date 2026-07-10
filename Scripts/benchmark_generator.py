from __future__ import annotations

from pathlib import Path
import csv
from typing import Any

from PIL import Image
from openpyxl import load_workbook
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from image_processor import combine_vertical

ROOT = Path(__file__).resolve().parents[1]
DATABASE_XLSX = ROOT / "Database" / "products.xlsx"
DATABASE_CSV = ROOT / "Database" / "products.csv"
OUTPUT = ROOT / "Output" / "Residential_ESS_Dimension_Benchmark.pptx"
TEMP_DIR = ROOT / "Output" / "_temp"

COLORS = {
    "navy": "0B1F33",
    "blue": "2563EB",
    "red": "DC2626",
    "green": "16A34A",
    "text": "172B4D",
    "muted": "607080",
    "line": "D9E2EC",
    "light": "F5F7FA",
    "white": "FFFFFF",
}


def rgb(hex_value: str) -> RGBColor:
    return RGBColor.from_string(hex_value.replace("#", ""))


def _normalize_record(record: dict[str, Any]) -> dict[str, Any]:
    width = int(float(record["Width_mm"]))
    depth = int(float(record["Depth_mm"]))
    height = int(float(record["Height_mm"]))
    return {
        "brand": str(record["Brand"]),
        "product": str(record["Product"]),
        "configuration": str(record["Configuration"]),
        "width_mm": width,
        "depth_mm": depth,
        "height_mm": height,
        "footprint_m2": width * depth / 1_000_000,
        "volume_m3": width * depth * height / 1_000_000_000,
        "asset_primary": str(record["Asset_primary"]),
        "asset_secondary": str(record.get("Asset_secondary") or ""),
        "notes": str(record.get("Notes") or ""),
    }


def load_products() -> list[dict[str, Any]]:
    products: list[dict[str, Any]] = []
    if DATABASE_XLSX.exists():
        workbook = load_workbook(DATABASE_XLSX, data_only=False)
        sheet = workbook["Products"]
        headers = [cell.value for cell in sheet[2]]
        for row in sheet.iter_rows(min_row=3, values_only=True):
            if not row[0]:
                continue
            products.append(_normalize_record(dict(zip(headers, row))))
    elif DATABASE_CSV.exists():
        with DATABASE_CSV.open("r", encoding="utf-8-sig", newline="") as handle:
            for record in csv.DictReader(handle):
                products.append(_normalize_record(record))
    else:
        raise FileNotFoundError("No product database found.")
    return products


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: float = 12,
    bold: bool = False,
    color: str = "172B4D",
    align=PP_ALIGN.LEFT,
    font: str = "Aptos",
    valign=MSO_ANCHOR.MIDDLE,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = Inches(0.02)
    frame.margin_top = frame.margin_bottom = Inches(0)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    run = paragraph.runs[0]
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: str = "FFFFFF",
    line: str = "D9E2EC",
    radius: bool = True,
    line_width: float = 0.8,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(line_width)
    return shape


def add_image_contain(slide, image_path: Path, center_x: float, bottom_y: float, height: float) -> float:
    with Image.open(image_path) as image:
        img_w, img_h = image.size
    width = height * img_w / img_h
    slide.shapes.add_picture(
        str(image_path),
        Inches(center_x - width / 2),
        Inches(bottom_y - height),
        width=Inches(width),
        height=Inches(height),
    )
    return width


def add_axis(slide, x: float, y: float, scale: float = 1.0):
    x_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.38 * scale), Inches(0.025)
    )
    x_line.fill.solid()
    x_line.fill.fore_color.rgb = rgb(COLORS["red"])
    x_line.line.fill.background()
    add_text(slide, "X", x + 0.39 * scale, y - 0.08, 0.16, 0.16, 8, True, COLORS["red"])

    z_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y - 0.36 * scale), Inches(0.025), Inches(0.36 * scale)
    )
    z_line.fill.solid()
    z_line.fill.fore_color.rgb = rgb(COLORS["blue"])
    z_line.line.fill.background()
    add_text(slide, "Z", x - 0.07, y - 0.53 * scale, 0.16, 0.16, 8, True, COLORS["blue"])

    y_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x + 0.02), Inches(y - 0.02), Inches(0.28 * scale), Inches(0.025)
    )
    y_line.rotation = -35
    y_line.fill.solid()
    y_line.fill.fore_color.rgb = rgb(COLORS["green"])
    y_line.line.fill.background()
    add_text(slide, "Y", x + 0.26 * scale, y - 0.28 * scale, 0.16, 0.16, 8, True, COLORS["green"])


def add_human(slide, x: float, baseline_y: float, height: float):
    head_y = baseline_y - height
    head = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.26), Inches(head_y), Inches(0.34), Inches(0.34))
    head.fill.solid()
    head.fill.fore_color.rgb = rgb("AAB7C4")
    head.line.fill.background()

    body_y = head_y + 0.38
    body_h = height - 1.45
    body = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.22), Inches(body_y), Inches(0.42), Inches(body_h)
    )
    body.fill.solid()
    body.fill.fore_color.rgb = rgb("AAB7C4")
    body.line.fill.background()

    for dx in (0.15, 0.49):
        leg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x + dx), Inches(baseline_y - 0.95), Inches(0.14), Inches(0.95)
        )
        leg.fill.solid()
        leg.fill.fore_color.rgb = rgb("AAB7C4")
        leg.line.fill.background()

    for dx in (0.04, 0.65):
        arm = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x + dx), Inches(body_y + 0.38), Inches(0.12), Inches(1.35)
        )
        arm.fill.solid()
        arm.fill.fore_color.rgb = rgb("AAB7C4")
        arm.line.fill.background()

    add_text(slide, "175 cm", x - 0.02, baseline_y + 0.05, 0.95, 0.22, 9, True, COLORS["muted"], PP_ALIGN.CENTER)


def prepare_asset(product: dict[str, Any]) -> Path:
    primary = ROOT / product["asset_primary"]
    secondary = product["asset_secondary"].strip()
    if secondary:
        TEMP_DIR.mkdir(parents=True, exist_ok=True)
        combined = TEMP_DIR / f"{product['brand']}_combined.png"
        return combine_vertical(ROOT / secondary, primary, combined)
    return primary


def add_bar_card(slide, title: str, subtitle: str, data, x: float, y: float, w: float, h: float):
    add_rect(slide, x, y, w, h, fill="FFFFFF", line=COLORS["line"], radius=True)
    add_text(slide, title, x + 0.18, y + 0.12, w - 0.36, 0.28, 14, True, COLORS["navy"])
    add_text(slide, subtitle, x + 0.18, y + 0.39, w - 0.36, 0.18, 8.5, False, COLORS["muted"])
    max_value = max(value for _, value, _ in data)
    row_h = (h - 0.75) / len(data)
    for idx, (name, value, label) in enumerate(data):
        row_y = y + 0.68 + idx * row_h
        add_text(slide, f"{idx + 1}", x + 0.12, row_y, 0.20, row_h * 0.75, 9, True, COLORS["muted"], PP_ALIGN.CENTER)
        add_text(slide, name, x + 0.34, row_y, 0.86, row_h * 0.75, 9, False, COLORS["text"])
        bar_x = x + 1.20
        full_bar_w = w - 2.05
        bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(bar_x),
            Inches(row_y + 0.05),
            Inches(full_bar_w),
            Inches(0.12),
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = rgb("EDF2F7")
        bg.line.fill.background()
        bar_w = max(0.16, full_bar_w * value / max_value)
        bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(bar_x),
            Inches(row_y + 0.05),
            Inches(bar_w),
            Inches(0.12),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = rgb(COLORS["blue"])
        bar.line.fill.background()
        add_text(slide, label, x + w - 0.75, row_y - 0.01, 0.60, row_h * 0.75, 8.5, True, COLORS["text"], PP_ALIGN.RIGHT)


def build_presentation(products: list[dict[str, Any]]) -> Path:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]

    # Slide 1: true-height product comparison
    slide = presentation.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(COLORS["white"])
    add_text(slide, "Residential ESS Dimension Benchmark", 0.48, 0.22, 8.5, 0.42, 24, True, COLORS["navy"])
    add_text(
        slide,
        "10 kW three-phase inverter + 1 battery pack · official product images · true-height comparison",
        0.48,
        0.66,
        10.8,
        0.26,
        10.5,
        False,
        COLORS["muted"],
    )
    add_text(slide, "V1.0", 11.85, 0.28, 0.95, 0.28, 10, True, COLORS["blue"], PP_ALIGN.RIGHT)

    baseline_y = 5.35
    baseline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.30), Inches(baseline_y), Inches(12.7), Inches(0.018)
    )
    baseline.fill.solid()
    baseline.fill.fore_color.rgb = rgb(COLORS["line"])
    baseline.line.fill.background()

    add_human(slide, 0.25, baseline_y, 4.25)

    start_x = 1.40
    col_w = 2.35
    max_img_h = 2.65
    max_height_mm = max(product["height_mm"] for product in products)
    height_scale = max_img_h / max_height_mm

    for idx, product in enumerate(products):
        col_x = start_x + idx * col_w
        center_x = col_x + col_w / 2
        add_text(slide, product["brand"], col_x, 1.03, col_w, 0.28, 12, True, COLORS["navy"], PP_ALIGN.CENTER)
        add_text(slide, product["product"], col_x + 0.08, 1.31, col_w - 0.16, 0.36, 8.5, False, COLORS["muted"], PP_ALIGN.CENTER)
        add_rect(slide, col_x + 0.16, 1.72, col_w - 0.32, 3.54, fill="FAFBFC", line="E6ECF2", radius=True)
        image_path = prepare_asset(product)
        image_h = product["height_mm"] * height_scale
        add_image_contain(slide, image_path, center_x, baseline_y - 0.10, image_h)
        add_axis(slide, col_x + 0.28, baseline_y - 0.14, 0.72)
        dimensions = f"{product['width_mm']} × {product['depth_mm']} × {product['height_mm']} mm"
        add_text(slide, dimensions, col_x + 0.12, baseline_y - 0.02, col_w - 0.24, 0.28, 9, True, COLORS["text"], PP_ALIGN.CENTER)

    # Top-view footprint panel
    add_rect(slide, 0.48, 5.72, 12.36, 1.40, fill="F5F7FA", line="E1E8EF", radius=True)
    add_text(slide, "Top View Footprint", 0.68, 5.84, 2.1, 0.24, 12, True, COLORS["navy"])
    add_text(slide, "Width × Depth, true relative proportion", 2.55, 5.84, 3.2, 0.24, 8.8, False, COLORS["muted"])
    fp_scale = min(
        1.35 / max(product["width_mm"] for product in products),
        0.48 / max(product["depth_mm"] for product in products),
    )
    fp_start_x = 1.38
    for idx, product in enumerate(products):
        center_x = fp_start_x + idx * col_w + col_w / 2
        footprint_w = product["width_mm"] * fp_scale
        footprint_h = product["depth_mm"] * fp_scale
        rect_x = center_x - footprint_w / 2
        rect_y = 6.22 + (0.50 - footprint_h) / 2
        add_rect(slide, rect_x, rect_y, footprint_w, footprint_h, fill="DCE9F7", line="86A9CC", radius=False)
        add_text(slide, product["brand"], center_x - 0.70, 6.71, 1.40, 0.20, 8.5, True, COLORS["text"], PP_ALIGN.CENTER)
        add_text(slide, f"{product['footprint_m2']:.3f} m²", center_x - 0.70, 6.91, 1.40, 0.18, 8, False, COLORS["muted"], PP_ALIGN.CENTER)

    # Slide 2: ranking dashboard
    slide = presentation.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(COLORS["white"])
    add_text(slide, "Compactness Ranking & Database Snapshot", 0.48, 0.22, 9.5, 0.42, 23, True, COLORS["navy"])
    add_text(
        slide,
        "Lower values rank better. Footprint and volume are calculated automatically from X / Y / Z dimensions.",
        0.48,
        0.66,
        11.6,
        0.26,
        10.2,
        False,
        COLORS["muted"],
    )

    height_data = [
        (p["brand"], p["height_mm"], f"{p['height_mm']} mm")
        for p in sorted(products, key=lambda item: (item["height_mm"], item["brand"]))
    ]
    footprint_data = [
        (p["brand"], p["footprint_m2"], f"{p['footprint_m2']:.3f} m²")
        for p in sorted(products, key=lambda item: (item["footprint_m2"], item["brand"]))
    ]
    volume_data = [
        (p["brand"], p["volume_m3"], f"{p['volume_m3']:.3f} m³")
        for p in sorted(products, key=lambda item: (item["volume_m3"], item["brand"]))
    ]

    add_bar_card(slide, "Height Ranking (Z)", "Lower is better", height_data, 0.48, 1.06, 4.00, 2.45)
    add_bar_card(slide, "Footprint Ranking (X×Y)", "Lower is better", footprint_data, 4.67, 1.06, 4.00, 2.45)
    add_bar_card(slide, "Envelope Volume", "Lower is better", volume_data, 8.86, 1.06, 4.00, 2.45)

    # Data table
    table_x, table_y, table_w, table_h = 0.48, 3.78, 9.25, 2.95
    add_rect(slide, table_x, table_y, table_w, table_h, fill="FFFFFF", line=COLORS["line"], radius=True)
    add_text(slide, "Dimension Database Snapshot", table_x + 0.20, table_y + 0.12, 4.0, 0.28, 13, True, COLORS["navy"])

    headers = ["Brand", "Width", "Depth", "Height", "Footprint", "Volume"]
    column_widths = [1.65, 1.25, 1.25, 1.25, 1.55, 1.55]
    row_h = 0.38
    table_start_x = table_x + 0.18
    table_start_y = table_y + 0.56
    current_x = table_start_x
    for header, column_w in zip(headers, column_widths):
        cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(current_x), Inches(table_start_y), Inches(column_w), Inches(row_h))
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb("DDEBF7")
        cell.line.color.rgb = rgb(COLORS["line"])
        add_text(slide, header, current_x + 0.02, table_start_y + 0.01, column_w - 0.04, row_h - 0.02, 8.8, True, COLORS["navy"], PP_ALIGN.CENTER)
        current_x += column_w

    for row_index, product in enumerate(products, start=1):
        values = [
            product["brand"],
            f"{product['width_mm']} mm",
            f"{product['depth_mm']} mm",
            f"{product['height_mm']} mm",
            f"{product['footprint_m2']:.3f} m²",
            f"{product['volume_m3']:.3f} m³",
        ]
        current_x = table_start_x
        current_y = table_start_y + row_index * row_h
        for value, column_w in zip(values, column_widths):
            cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(current_x), Inches(current_y), Inches(column_w), Inches(row_h))
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb("FFFFFF" if row_index % 2 else "F8FAFC")
            cell.line.color.rgb = rgb(COLORS["line"])
            add_text(slide, value, current_x + 0.02, current_y + 0.01, column_w - 0.04, row_h - 0.02, 8.6, False, COLORS["text"], PP_ALIGN.CENTER)
            current_x += column_w

    # Insights panel
    add_rect(slide, 9.98, 3.78, 2.88, 2.95, fill="F5F7FA", line=COLORS["line"], radius=True)
    add_text(slide, "Design Insights", 10.18, 3.92, 2.45, 0.30, 13, True, COLORS["navy"])
    shortest = min(products, key=lambda item: item["height_mm"])
    smallest_fp = min(products, key=lambda item: item["footprint_m2"])
    smallest_vol = min(products, key=lambda item: item["volume_m3"])
    deepest = max(products, key=lambda item: item["depth_mm"])
    insights = [
        f"Shortest system: {shortest['brand']} ({shortest['height_mm']} mm)",
        f"Smallest footprint: {smallest_fp['brand']} ({smallest_fp['footprint_m2']:.3f} m²)",
        f"Smallest envelope volume: {smallest_vol['brand']} ({smallest_vol['volume_m3']:.3f} m³)",
        f"Largest depth: {deepest['brand']} ({deepest['depth_mm']} mm)",
        "Official product images only; no AI redraw.",
    ]
    insight_y = 4.36
    for item in insights:
        add_text(slide, "•", 10.18, insight_y, 0.18, 0.34, 10, True, COLORS["blue"], PP_ALIGN.CENTER)
        add_text(slide, item, 10.38, insight_y, 2.22, 0.42, 9.3, False, COLORS["text"], PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP)
        insight_y += 0.44

    add_text(slide, "X = Width (red)  |  Y = Depth (green)  |  Z = Height (blue)", 0.52, 7.08, 7.0, 0.18, 8.5, False, COLORS["muted"])

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    output = build_presentation(load_products())
    print(f"Generated: {output}")
